import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

from entidades.Conta import Conta


def importar_dados_excel(caminho_arquivo):
    """
    Importa os dados do arquivo Excel para um DataFrame.
    """
    df = pd.read_excel(caminho_arquivo)
    return df


def registrar_valores(df_filtrado, balanco, encontrar_conta_por_codigo):
    """
    Registra os valores no balanço patrimonial com base nos dados filtrados.
    """
    for index, row in df_filtrado.iterrows():
        codigo = row['Cod']
        valor = row['Valor']
        conta = encontrar_conta_por_codigo(balanco.ativo.contas, codigo) or encontrar_conta_por_codigo(
            balanco.passivo.contas, codigo)

        if conta:
            if conta.subcontas:
                subconta = encontrar_conta_por_codigo(conta.subcontas, codigo)
                if subconta:
                    subconta.atualizar_valor(valor)
                else:
                    # Adiciona o valor à subconta 'Não Especificado' se a subconta específica não existir
                    subconta_nao_especificado = encontrar_conta_por_codigo(conta.subcontas, 'Não Especificado')
                    if subconta_nao_especificado:
                        subconta_nao_especificado.atualizar_valor(valor)
                    else:
                        # Caso a subconta 'Não Especificado' não exista, um erro é lançado
                        raise ValueError(f"Subconta 'Não Especificado' não encontrada para a conta {conta.nome}.")
            else:
                # Atualiza a conta principal diretamente
                conta.atualizar_valor(valor)
    return balanco


def filtrar_dados(df, mes, cnpj):
    """
    Filtra os dados com base no mês e no CNPJ.
    """
    # Convertendo a coluna 'Mes' para o formato YYYY-MM
    df['Mes'] = df['Mes'].dt.strftime('%Y-%m')

    # Converte CNPJ para string para garantir consistência na comparação
    df['CNPJ'] = df['CNPJ'].astype(str)

    # Aplicando o filtro
    df_filtrado = df[(df['Mes'] == mes) & (df['CNPJ'] == str(cnpj))]
    return df_filtrado


def construir_contas(dados, codigo=""):
    if not isinstance(dados, dict):
        raise TypeError(f"Esperado um dicionário, mas encontrou {type(dados)} para o código {codigo}")

    contas = []
    for sub_codigo, sub_dados in dados.items():
        if not isinstance(sub_dados, dict):
            raise TypeError(
                f"Esperado um dicionário em sub_dados, mas encontrou {type(sub_dados)} para o código {sub_codigo}")

        conta = Conta(sub_dados["nome"], sub_dados["valor"], codigo=sub_codigo)

        if "subcategorias" in sub_dados:
            if not isinstance(sub_dados["subcategorias"], dict):
                raise TypeError(
                    f"Esperado um dicionário em subcategorias, mas encontrou {type(sub_dados['subcategorias'])} para o código {sub_codigo}")
            subcontas = construir_contas(sub_dados["subcategorias"], sub_codigo)
            for subconta in subcontas:
                conta.adicionar_subconta(subconta)
        contas.append(conta)
    return contas


def encontrar_conta_por_nome(grupo_contas, nome):
    """
    Encontra uma conta pelo nome dentro da hierarquia de contas.
    """
    for conta in grupo_contas:
        if conta.nome == nome:
            return conta
        for subconta in conta.subcontas:
            resultado = encontrar_conta_por_nome([subconta], nome)
            if resultado:
                return resultado
    return None


def encontrar_conta_por_codigo(grupo_contas, codigo):
    """
    Encontra uma conta pelo código dentro da hierarquia de contas.
    """
    for conta in grupo_contas:
        if isinstance(conta, Conta) and conta.codigo == codigo:
            return conta
        if isinstance(conta, Conta):  # Verifica se é um objeto Conta antes de acessar subcontas
            for subconta in conta.subcontas:
                resultado = encontrar_conta_por_codigo([subconta], codigo)
                if resultado:
                    return resultado
    return None


def montar_df_plano_contas(balanco, df_filtrado):
    """
    Constrói um DataFrame com a estrutura completa do plano de contas,
    incluindo informações adicionais de competência e empresa.
    """
    data = []

    def adicionar_contas(contas, nivel=0):
        for conta in contas:
            linha = {
                "Código": conta.codigo,
                "Nome": conta.nome,
                "Nível": nivel,
                "Valor": conta.calcular_valor_total(),
                # As informações de competência e empresa são mantidas para cada linha
                "Mes": df_filtrado["Mes"].iloc[0],  # Assumimos que todas as linhas de df_filtrado têm o mesmo mês
                "Empresa": df_filtrado["Empresa"].iloc[0],
                # Assumimos que todas as linhas de df_filtrado têm a mesma empresa
            }
            data.append(linha)
            if conta.subcontas:
                adicionar_contas(conta.subcontas, nivel + 1)

    adicionar_contas(balanco.ativo.contas, nivel=0)
    adicionar_contas(balanco.passivo.contas, nivel=0)

    df = pd.DataFrame(data)
    return df


def calcular_variacoes(balanco_1, balanco_2, nome_plano_1, nome_plano_2):
    """
    Calcula as variações absolutas e percentuais entre dois balanços patrimoniais,
    tratando adequadamente os itens 'Não Especificado' dentro de seus contextos.
    """
    df_filtrado_1 = st.session_state["df_filtrado"]
    df_filtrado_2 = st.session_state["df_filtrado"]

    # Convertendo os balanços em DataFrames com o df_filtrado
    df_1 = montar_df_plano_contas(balanco_1, df_filtrado_1)
    df_2 = montar_df_plano_contas(balanco_2, df_filtrado_2)

    # Criar uma chave de contexto única para identificar adequadamente contas 'Não Especificado'
    df_1["Chave Contexto"] = df_1.apply(
        lambda row: f"{row['Código']}|{row['Nome']}|{obter_hierarquia_contexto(row['Código'])}", axis=1
    )
    df_2["Chave Contexto"] = df_2.apply(
        lambda row: f"{row['Código']}|{row['Nome']}|{obter_hierarquia_contexto(row['Código'])}", axis=1
    )

    # Realizando o merge dos dois DataFrames com base na chave composta
    df_merged = pd.merge(df_1, df_2, on="Chave Contexto", suffixes=("_1", "_2"))

    # Calculando a variação absoluta e percentual entre os valores dos dois planos de contas
    df_merged["Variação Absoluta"] = df_merged["Valor_2"] - df_merged["Valor_1"]
    df_merged["Variação Percentual (%)"] = (
            (df_merged["Variação Absoluta"] / df_merged["Valor_1"].replace(0, float("nan"))) * 100
    ).round(2)

    # Filtrando apenas as colunas relevantes para exibição
    df_variacoes = df_merged[
        ["Código_1", "Nome_1", "Valor_1", "Valor_2", "Variação Absoluta", "Variação Percentual (%)"]
    ]

    # Renomeando as colunas para melhorar a leitura e incluir os nomes dos planos
    df_variacoes.rename(
        columns={
            "Código_1": "Código",
            "Nome_1": "Nome",
            "Valor_1": f"Valor ({nome_plano_1})",
            "Valor_2": f"Valor ({nome_plano_2})",
        },
        inplace=True,
    )

    return df_variacoes


def obter_hierarquia_contexto(codigo):
    """
    Retorna um identificador de contexto hierárquico baseado no código da conta.
    Esse identificador ajuda a distinguir contas 'Não Especificado' em contextos diferentes.
    """
    # Quebra o código pelo separador para obter uma representação hierárquica do contexto
    segmentos = codigo.split(".")
    return ".".join(segmentos[:-1]) if len(segmentos) > 1 else codigo


def consolidar_nao_especificado(df, nome_plano_1, nome_plano_2):
    """
    Consolida valores para contas 'Não Especificado' dentro do mesmo contexto hierárquico,
    evitando duplicações indesejadas.

    :param nome_plano_1: Nome do primeiro plano de contas.
    :param nome_plano_2: Nome do segundo plano de contas.
    :param df: DataFrame com as variações calculadas.
    :return: DataFrame ajustado com contas 'Não Especificado' consolidadas por contexto.
    """
    df_nao_especificado = df[df["Nome"].str.contains("Não Especificado")]

    # Agrupa as contas 'Não Especificado' pelo contexto identificador e soma os valores
    df_consolidado = df_nao_especificado.groupby(["Código", "Nome"]).agg({
        f"Valor ({nome_plano_1})": 'sum',
        f"Valor ({nome_plano_2})": 'sum',
        "Variação Absoluta": 'sum',
        "Variação Percentual (%)": 'mean'  # Ajuste se necessário, dependendo do contexto
    }).reset_index()

    # Substitui as contas 'Não Especificado' no DataFrame original com as consolidadas
    df = df[~df["Nome"].str.contains("Não Especificado")]
    df = pd.concat([df, df_consolidado], ignore_index=True)

    return df


def filtrar_contas_ultimos_12_meses(df: pd.DataFrame) -> pd.DataFrame:
    """
    Filtra as contas que tiveram valores lançados nos últimos 12 meses e limita a profundidade a 3 níveis.

    :param df: DataFrame contendo as contas e valores.
    :return: DataFrame filtrado.
    """
    # Converte a coluna 'Mes' para o formato datetime padrão se necessário
    df = converter_para_datetime(df, 'Mes', '%Y-%m')

    # Limitar ao nível 3 da hierarquia (contas com códigos até 3 níveis)
    df['Nivel'] = df['Cod'].apply(lambda x: len(x.split('.')))
    df_filtrado = df[df['Nivel'] <= 3]

    # Filtrar contas com valor diferente de 0
    df_filtrado = df_filtrado[df_filtrado['Valor'] != 0]

    # Filtrar pelos últimos 12 meses
    ultimo_mes = df_filtrado['Mes'].max()
    df_filtrado_12_meses = df_filtrado[df_filtrado['Mes'] >= (ultimo_mes - pd.DateOffset(months=12))]

    return df_filtrado_12_meses



def exportar_para_pptx(df: pd.DataFrame, caminho_pptx: str):
    """
    Exporta as contas filtradas para um arquivo PowerPoint, criando uma tabela no slide.

    :param df: DataFrame filtrado com as contas.
    :param caminho_pptx: Caminho para salvar o arquivo PowerPoint.
    """
    # Criar apresentação
    prs = Presentation()

    # Adicionar slide com layout de título e conteúdo
    slide_layout = prs.slide_layouts[5]  # layout vazio
    slide = prs.slides.add_slide(slide_layout)

    # Adicionar título ao slide
    title = slide.shapes.title
    title.text = "Relatório de Contas (Últimos 12 Meses)"

    # Adicionar tabela ao slide
    rows, cols = df.shape[0] + 1, df.shape[1]  # +1 para o cabeçalho
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table

    # Adicionar cabeçalho da tabela
    for i, col_name in enumerate(df.columns):
        table.cell(0, i).text = col_name

    # Preencher os dados da tabela
    for i, row in enumerate(df.itertuples(), start=1):
        for j, value in enumerate(row[1:], start=0):
            table.cell(i, j).text = str(value)

    # Salvar a apresentação
    prs.save(caminho_pptx)
    st.success(f"Apresentação salva em: {caminho_pptx}")


def converter_para_datetime(df, coluna_data='Mes', formato='%Y-%m'):
    """
    Converte a coluna de data para datetime com um formato padrão.

    :param df: DataFrame com a coluna de data.
    :param coluna_data: Nome da coluna de data.
    :param formato: Formato da data a ser convertido.
    :return: DataFrame com a coluna de data convertida para datetime.
    """
    try:
        df[coluna_data] = pd.to_datetime(df[coluna_data], format=formato, errors='coerce')
    except Exception as e:
        raise ValueError(f"Erro ao converter a coluna '{coluna_data}' para datetime: {e}")
    return df
